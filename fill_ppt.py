from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
from datetime import datetime
import os
import sys


def run_fill_ppt(excel_path=None, output_dir="Output", output_filename=None, menu_type="A"):
    """Generate PPT from the provided Excel file using the template.

    `excel_path` must be provided by the user and point to an existing .xlsx file.
    PPT template is fixed at templates/ppt_temp.pptx.
    `menu_type` can be "A" or "B" to specify which meal type to generate.
    Returns the output path on success.
    """
    # 要求用户手动指定 Excel 文件路径（不再自动查找最新文件）
    if not excel_path:
        raise ValueError("请提供要处理的 Excel 文件路径：excel_path 参数不能为空")

    if not os.path.isfile(excel_path):
        raise FileNotFoundError(f"未找到指定的 Excel 文件: {excel_path}")

    print(f"[INFO] 使用指定 Excel 文件：{excel_path}")
    
    # 获取正确的模板路径（支持 EXE 和开发环境）
    # 如果是 PyInstaller 打包的 EXE，使用 sys._MEIPASS；否则使用当前目录
    if hasattr(sys, '_MEIPASS'):
        base_path = sys._MEIPASS
    else:
        base_path = os.path.dirname(os.path.abspath(__file__))
    
    ppt_path = os.path.join(base_path, "templates", "ppt_temp.pptx")
    
    # 验证模板文件是否存在
    if not os.path.isfile(ppt_path):
        raise FileNotFoundError(f"未找到 PPT 模板文件: {ppt_path}\n基础路径: {base_path}")

    # 尝试读取Excel文件
    try:
        wb = load_workbook(excel_path)
        
        # 查找包含"二4"的工作表
        target_sheet_name = next((name for name in wb.sheetnames if "二4" in name), None)
        
        if not target_sheet_name:
            raise ValueError(f"未找到包含 '二4' 的 sheet\n可用的工作表: {', '.join(wb.sheetnames)}")
        
        ws = wb[target_sheet_name]
        
    except Exception as e:
        raise ValueError(f"读取Excel文件失败: {str(e)}")

    # 获取日期（第二行 B-F 列，即 row 2, columns B-F）
    formatted_dates = []
    for col_idx in range(2, 7):  # B-F: columns 2-6
        cell_value = ws.cell(row=2, column=col_idx).value
        if cell_value:
            if isinstance(cell_value, datetime):
                dt = cell_value
            else:
                try:
                    dt = datetime.strptime(str(cell_value), "%Y-%m-%d")
                except:
                    dt = cell_value
            if isinstance(dt, datetime):
                formatted_dates.append(f"{dt.month}月{dt.day}日")
            else:
                formatted_dates.append(str(cell_value))
        else:
            formatted_dates.append("")

    # 获取每天选择指定餐类型的名单
    meal_choices = {}
    for i, col_idx in enumerate(range(2, 7)):  # B-F: columns 2-6
        day_students = []
        for row_idx in range(3, ws.max_row + 1):  # 从第3行开始
            name_cell = ws.cell(row=row_idx, column=1).value
            meal_cell = ws.cell(row=row_idx, column=col_idx).value
            if meal_cell == menu_type and name_cell:
                day_students.append(str(name_cell))
        meal_choices[i] = day_students

    # 获取"共计 x 份"数据（A列为"{menu_type}餐合计"）
    meal_counts = []
    count_label = f"{menu_type}餐合计"
    for col_idx in range(2, 7):  # B-F: columns 2-6
        count = 0
        for row_idx in range(3, ws.max_row + 1):
            label_cell = ws.cell(row=row_idx, column=1).value
            if label_cell and str(label_cell) == count_label:
                count_cell = ws.cell(row=row_idx, column=col_idx).value
                if count_cell:
                    count = int(count_cell) if isinstance(count_cell, (int, float)) else 0
                break
        meal_counts.append(count)

    # 校验每天统计人数与"{menu_type}餐合计"是否一致
    for i in range(5):
        counted = len(meal_choices[i])
        expected = meal_counts[i]
        if counted == expected:
            print(f"[确认] 第{i+1}天 {menu_type}餐人数正确：{counted}人")
        else:
            print(f"[警告] 第{i+1}天 {menu_type}餐人数不一致！统计为 {counted} 人，{menu_type}餐合计为 {expected} 人")

    # 加载 PPT 模板
    prs = Presentation(ppt_path)

    # 设置表格字体样式
    def format_table_cell(cell, name):
        cell.text = name
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = 'SimSun'
                run.font.size = Pt(32)
                run.font.bold = True

    # 设置日期样式
    def format_date_text(shape, date_str):
        shape.text = date_str
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(54)
                run.font.color.rgb = RGBColor(255, 102, 0)  # 橙色
                run.font.shadow = True
                run.font.bold = True

    # 设置“共计_x__份”样式
    def format_total_text(shape, count):
        shape.text = f"共计_{count}__份"
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = '字魂95号-手刻宋'
                run.font.size = Pt(54)
                run.font.color.rgb = RGBColor(0, 102, 255)  # 蓝色
                run.font.reflection = True

    # 填充第3至第7页幻灯片
    for i in range(5):
        slide = prs.slides[i + 2]
        shapes = slide.shapes

        # 插入日期
        for shape in shapes:
            if shape.has_text_frame and "月" in shape.text and "日" in shape.text:
                format_date_text(shape, formatted_dates[i])
                break

        # 插入“共计_x__份”
        for shape in shapes:
            if shape.has_text_frame and "共计" in shape.text and "份" in shape.text:
                format_total_text(shape, meal_counts[i])
                break

        # 插入学生名单（表格）
        table = None
        for shape in shapes:
            if shape.has_table:
                table = shape.table
                break

        if table:
            names = meal_choices[i]
            row_idx = 0
            col_idx = 0
            for name in names:
                if row_idx >= len(table.rows):
                    table.add_row()
                cell = table.cell(row_idx, col_idx)
                format_table_cell(cell, name)
                col_idx += 1
                if col_idx == 7:
                    col_idx = 0
                    row_idx += 1

    # 保存新 PPT 文件
    if output_filename is None:
        # 计算周数（以 9月1日为第一周）
        first_date_str = ws.cell(row=2, column=2).value  # 取第一天的日期
        if isinstance(first_date_str, datetime):
            first_date = first_date_str
        else:
            try:
                first_date = datetime.strptime(str(first_date_str), "%Y-%m-%d")
            except:
                first_date = datetime.now()
        
        semester_start = datetime(2025, 9, 1)
        week_number = ((first_date - semester_start).days // 7) + 1

        # 生成文件名
        output_filename = f"文明用餐二上（第{week_number}周）.pptx"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, output_filename)
    prs.save(output_path)
    print(f"新PPT已保存至: {output_path}")
    return output_path


if __name__ == '__main__':
    # 简单的命令行接口：
    # 用法: python fill_ppt.py <excel_path> [output_dir]
    if len(sys.argv) < 2:
        print("用法: python fill_ppt.py <excel_path> [output_dir]")
        print("示例: python fill_ppt.py Input/roster.xlsx Output")
        sys.exit(1)

    excel = sys.argv[1]
    out_dir = sys.argv[2] if len(sys.argv) > 2 else "Output"
    run_fill_ppt(excel_path=excel, output_dir=out_dir)